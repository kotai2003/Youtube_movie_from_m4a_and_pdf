"""
Step 1 — Slide Extraction

Extracts text content and images from presentation files (PPTX or PDF).
Falls back to EasyOCR for image-based slides that contain no extractable text.

Supported formats:
    - **.pptx** — via python-pptx (text) + PowerPoint COM / LibreOffice (images)
    - **.pdf** — via PyMuPDF (text + rasterisation)
"""

import os
import json
import sys
from pathlib import Path
from PIL import Image


# --- OCR ---

_ocr_reader = None

def _get_ocr_reader(languages: list[str] = None):
    """EasyOCRリーダーをシングルトンで取得（初回のみモデルダウンロード）"""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        if languages is None:
            languages = ["ja", "en"]
        print(f"  OCRモデルをロード中（{languages}）...")
        _ocr_reader = easyocr.Reader(languages, gpu=True)
    return _ocr_reader


def _ocr_image(image_path: str, languages: list[str] = None) -> str:
    """画像からOCRでテキストを抽出"""
    reader = _get_ocr_reader(languages)
    results = reader.readtext(image_path, detail=0, paragraph=True)
    return "\n".join(results)


def _is_text_empty(text: str) -> bool:
    """テキストが実質的に空かどうか判定"""
    cleaned = text.strip()
    if not cleaned:
        return True
    # 「スライド N」だけの場合も空とみなす
    if cleaned.replace(" ", "").replace("　", "").startswith("スライド"):
        return True
    # 数文字しかない場合
    if len(cleaned) < 5:
        return True
    return False


# --- PPTX ---

def extract_from_pptx(pptx_path: str, output_dir: str, ocr_languages: list[str] = None) -> list[dict]:
    """Extract text and images from a PPTX file.

    Parameters
    ----------
    pptx_path : str
        Path to the ``.pptx`` file.
    output_dir : str
        Directory where ``slide_images/`` will be created.
    ocr_languages : list[str], optional
        Languages for EasyOCR fallback (default ``["ja", "en"]``).

    Returns
    -------
    list[dict]
        Per-slide dicts with keys ``slide_number``, ``title``,
        ``full_text``, ``image_path``, and optionally ``ocr_used``.
    """
    from pptx import Presentation

    slides_info = []
    prs = Presentation(pptx_path)

    img_dir = os.path.join(output_dir, "slide_images")
    os.makedirs(img_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            if shape.shape_type is not None and hasattr(shape, "text"):
                if shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0:
                        title = shape.text.strip()

        if not title and texts:
            title = texts[0]

        slides_info.append({
            "slide_number": i,
            "title": title,
            "full_text": "\n".join(texts),
            "image_path": None
        })

    # PPTX → 画像変換
    try:
        _pptx_to_images_windows(pptx_path, img_dir, slides_info)
    except Exception as e:
        print(f"  [INFO] PowerPoint COM変換失敗: {e}")
        try:
            _pptx_to_images_libreoffice(pptx_path, img_dir, slides_info)
        except Exception as e2:
            print(f"  [WARN] LibreOffice変換も失敗: {e2}")
            print("  [WARN] スライドをPDFで保存し、PDFモードで再実行してください。")

    # テキストが空のスライドにはOCRを適用
    ocr_count = 0
    for s in slides_info:
        if _is_text_empty(s["full_text"]) and s.get("image_path") and os.path.exists(s["image_path"]):
            ocr_text = _ocr_image(s["image_path"], ocr_languages)
            if ocr_text.strip():
                s["full_text"] = ocr_text
                if _is_text_empty(s["title"]):
                    # OCRテキストの最初の行をタイトルに
                    first_line = ocr_text.split("\n")[0].strip()
                    s["title"] = first_line[:60] if first_line else f"スライド {s['slide_number']}"
                s["ocr_used"] = True
                ocr_count += 1

    if ocr_count > 0:
        print(f"  → {ocr_count} 枚のスライドでOCRを使用しました")

    return slides_info


def _pptx_to_images_windows(pptx_path: str, img_dir: str, slides_info: list):
    """Windows環境: PowerPoint COMオートメーションで画像化"""
    import comtypes.client

    pptx_abs = os.path.abspath(pptx_path)
    img_dir_abs = os.path.abspath(img_dir)

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    try:
        presentation = powerpoint.Presentations.Open(pptx_abs, WithWindow=False)

        for i, slide in enumerate(presentation.Slides, 1):
            img_path = os.path.join(img_dir_abs, f"slide_{i:03d}.png")
            slide.Export(img_path, "PNG", 1920, 1080)
            if i <= len(slides_info):
                slides_info[i - 1]["image_path"] = img_path
            print(f"  スライド {i} → {img_path}")

        presentation.Close()
    finally:
        powerpoint.Quit()


def _pptx_to_images_libreoffice(pptx_path: str, img_dir: str, slides_info: list):
    """LibreOfficeで画像化（フォールバック）"""
    import subprocess

    pptx_abs = os.path.abspath(pptx_path)

    result = subprocess.run([
        "soffice", "--headless", "--convert-to", "pdf",
        "--outdir", img_dir, pptx_abs
    ], capture_output=True, text=True, encoding="utf-8", errors="replace")

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice変換エラー: {result.stderr}")

    pdf_name = Path(pptx_path).stem + ".pdf"
    pdf_path = os.path.join(img_dir, pdf_name)

    _pdf_to_images(pdf_path, img_dir, slides_info)
    os.remove(pdf_path)


# --- PDF ---

def extract_from_pdf(pdf_path: str, output_dir: str, ocr_languages: list[str] = None) -> list[dict]:
    """Extract text and images from a PDF file.

    Parameters
    ----------
    pdf_path : str
        Path to the ``.pdf`` file.
    output_dir : str
        Directory where ``slide_images/`` will be created.
    ocr_languages : list[str], optional
        Languages for EasyOCR fallback (default ``["ja", "en"]``).

    Returns
    -------
    list[dict]
        Per-slide dicts (same schema as :func:`extract_from_pptx`).
    """
    import fitz  # PyMuPDF

    slides_info = []
    img_dir = os.path.join(output_dir, "slide_images")
    os.makedirs(img_dir, exist_ok=True)

    doc = fitz.open(pdf_path)

    for i, page in enumerate(doc, 1):
        # テキスト抽出
        text = page.get_text().strip()
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        title = lines[0] if lines else f"スライド {i}"

        # 画像化（高解像度）
        mat = fitz.Matrix(2.5, 2.5)
        pix = page.get_pixmap(matrix=mat)
        img_path = os.path.join(img_dir, f"slide_{i:03d}.png")
        pix.save(img_path)

        # 1920x1080にリサイズ
        img = Image.open(img_path)
        img = img.resize((1920, 1080), Image.LANCZOS)
        img.save(img_path)

        slides_info.append({
            "slide_number": i,
            "title": title,
            "full_text": text,
            "image_path": os.path.abspath(img_path)
        })

    doc.close()

    # テキストが空のスライドにはOCRを適用
    ocr_count = 0
    for s in slides_info:
        if _is_text_empty(s["full_text"]) and s.get("image_path") and os.path.exists(s["image_path"]):
            print(f"  スライド {s['slide_number']}: テキストなし → OCR実行中...")
            ocr_text = _ocr_image(s["image_path"], ocr_languages)
            if ocr_text.strip():
                s["full_text"] = ocr_text
                first_line = ocr_text.split("\n")[0].strip()
                s["title"] = first_line[:60] if first_line else f"スライド {s['slide_number']}"
                s["ocr_used"] = True
                ocr_count += 1

        img_path = s.get("image_path", "")
        title_display = s["title"][:40]
        print(f"  スライド {s['slide_number']}: 「{title_display}」 → {img_path}")

    if ocr_count > 0:
        print(f"\n  → {ocr_count} 枚のスライドでOCRを使用しました")

    return slides_info


def _pdf_to_images(pdf_path: str, img_dir: str, slides_info: list):
    """PDFを画像に変換（内部ヘルパー）"""
    import fitz

    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc, 1):
        mat = fitz.Matrix(2.5, 2.5)
        pix = page.get_pixmap(matrix=mat)
        img_path = os.path.join(img_dir, f"slide_{i:03d}.png")
        pix.save(img_path)

        img = Image.open(img_path)
        img = img.resize((1920, 1080), Image.LANCZOS)
        img.save(img_path)

        if i <= len(slides_info):
            slides_info[i - 1]["image_path"] = os.path.abspath(img_path)
    doc.close()


# --- メイン ---

def extract_slides(slides_file: str, output_dir: str, ocr_languages: list[str] = None) -> list[dict]:
    """High-level entry point for slide extraction.

    Dispatches to :func:`extract_from_pptx` or :func:`extract_from_pdf`
    based on the file extension, then saves ``slides_info.json``.

    Parameters
    ----------
    slides_file : str
        Path to ``.pptx`` or ``.pdf``.
    output_dir : str
        Output directory for images and JSON.
    ocr_languages : list[str], optional
        Languages for EasyOCR fallback.

    Returns
    -------
    list[dict]
        Per-slide information written to ``slides_info.json``.
    """
    ext = Path(slides_file).suffix.lower()

    print(f"\n{'='*50}")
    print(f"ステップ1: スライド抽出 ({ext})")
    print(f"{'='*50}")

    if ext == ".pptx":
        slides_info = extract_from_pptx(slides_file, output_dir, ocr_languages)
    elif ext == ".pdf":
        slides_info = extract_from_pdf(slides_file, output_dir, ocr_languages)
    else:
        raise ValueError(f"未対応のファイル形式: {ext}（.pptx または .pdf に対応）")

    # 結果保存
    info_path = os.path.join(output_dir, "slides_info.json")
    with open(info_path, "w", encoding="utf-8") as f:
        json.dump(slides_info, f, ensure_ascii=False, indent=2)

    print(f"\n  → {len(slides_info)} 枚のスライドを抽出しました")
    print(f"  → スライド情報: {info_path}")

    return slides_info


if __name__ == "__main__":
    import yaml
    with open("config.yaml", "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    os.makedirs(config["output_dir"], exist_ok=True)

    ocr_langs = config.get("ocr", {}).get("languages", ["ja", "en"])
    extract_slides(config["slides_file"], config["output_dir"], ocr_langs)
